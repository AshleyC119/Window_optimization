#!/usr/bin/env python3
"""Generate slides.pptx from project data."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

C1 = RGBColor(0x16, 0x21, 0x3E)  # dark accent
C2 = RGBColor(0x0F, 0x34, 0x60)
WH = RGBColor(0xFF, 0xFF, 0xFF)
DK = RGBColor(0x2D, 0x2D, 0x2D)
MU = RGBColor(0x77, 0x77, 0x77)
RD = RGBColor(0xC0, 0x39, 0x2B)
GY = RGBColor(0xF8, 0xF9, 0xFA)

def blank(): return prs.slides.add_slide(prs.slide_layouts[6])

def T(slide, txt, l, t, w, h, sz=13, b=False, c=DK, a=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    bx.text_frame.word_wrap = True
    p = bx.text_frame.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = c; p.alignment = a
    return bx.text_frame

def H(slide, txt, t=0.3, sz=27):
    return T(slide, txt, 0.5, t, 12.3, 0.6, sz, True, C1)

def Img(slide, path, l, t, w, h=None):
    if not os.path.exists(path): return
    if h: slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else: slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))

def Tbl(slide, data, l, t, cw):
    r,c = len(data), len(data[0])
    tbl = slide.shapes.add_table(r,c,Inches(l),Inches(t),Inches(sum(cw)),Inches(0.3*r)).table
    for i,w in enumerate(cw): tbl.columns[i].width = Inches(w)
    for ri,row in enumerate(data):
        for ci,val in enumerate(row):
            cl = tbl.cell(ri,ci); cl.text = str(val)
            for p in cl.text_frame.paragraphs:
                p.font.size = Pt(11)
                if ri==0: p.font.bold=True; p.font.color.rgb=WH
            if ri==0: cl.fill.solid(); cl.fill.fore_color.rgb=C1
            elif ri==r-1: cl.fill.solid(); cl.fill.fore_color.rgb=RGBColor(0xFE,0xF9,0xE7)

def Stat(slide, txt, l, t, w, sz=11):
    bx = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.35))
    bx.fill.solid(); bx.fill.fore_color.rgb = RGBColor(0xF0,0xF4,0xFF)
    bx.line.fill.background()
    bx.text_frame.word_wrap = True; bx.text_frame.margin_left=Pt(6)
    p = bx.text_frame.paragraphs[0]; p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = DK
    # Bold the numbers
    return bx

def Card(slide, title, body, l, t, w, h, border=False):
    s = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = GY
    if border: s.line.color.rgb = C2; s.line.width = Pt(2)
    else: s.line.fill.background()
    tf = s.text_frame; tf.word_wrap=True; tf.margin_left=Pt(8); tf.margin_right=Pt(8)
    p=tf.paragraphs[0]; p.text=title; p.font.size=Pt(12); p.font.bold=True; p.font.color.rgb=C1
    p2=tf.add_paragraph(); p2.text=body; p2.font.size=Pt(10); p2.font.color.rgb=DK
    return s

# ═══════════ SLIDE 1 ═══════════
s = blank()
T(s, '基于代理辅助演化算法的\n建筑电磁窗口多目标优化设计研究', 1.5, 1.5, 10.3, 2.2, 36, True, C1, PP_ALIGN.CENTER)
T(s, '从跨算法互证到 BO-LGBM 代理模型混合搜索框架', 2, 3.6, 9.3, 0.6, 18, False, C2, PP_ALIGN.CENTER)
T(s, '物理降噪 → 算法互证 → 代理模型 → 两阶段 Pipeline', 3, 4.8, 7.3, 0.5, 14, False, MU, PP_ALIGN.CENTER)
T(s, '汇报人：成雨鸣', 5, 6.2, 3.3, 0.5, 14, False, MU, PP_ALIGN.CENTER)

# ═══════════ SLIDE 2 ═══════════
s = blank()
H(s, '一、课题背景、数学模型与 KDE 网格积分等价性降噪')
T(s, '高频毫米波 O2I 穿透 Low-E 幕墙引入 20–30 dB 损耗。物理建模含：\n8 天线阵列均匀菲涅尔衍射信道 (LoS+NLoS)、人群 RWP 行为模型、\n四方向用户自阻塞模型。', 0.5, 1.2, 6.0, 1.2, 13)
T(s, 'θ = [x\u1d9c, z\u1d9c, L\u2093, L\u2d9e]\u1d40\nmin f₁ = Area    min f₂ = Outage', 1.5, 2.6, 4, 0.8, 14, False, C2, PP_ALIGN.CENTER)
T(s, 'MATLAB 300s 短时 RWP → σ≈4.5% 噪声 → 早停。\n重构：200×200×5 确定性网格积分 (KDE)。', 7, 1.2, 5.8, 0.7, 12)
Img(s, 'python_optimization/fig_exp1_convergence.png', 7, 1.9, 5.5, 2.1)
T(s, '20,000s 时偏差仅 0.67%。100 组随机配置：\nMAE=1.22%, R²=0.9941 —— 数学确证等价。', 7, 4.1, 5.5, 0.7, 12)
Img(s, 'python_optimization/fig_exp2_discrepancy.png', 7, 4.8, 5.5, 2.1)

# ═══════════ SLIDE 3 ═══════════
s = blank()
H(s, '二、多目标启发式算法体系与三大算法交叉互证')
T(s, '引入三种机制不同的启发性算法——NSGA-II (支配度)、AGE-MOEA (几何估计)、\nMOEA/D (标量化分解)——以排除单一算法局部最优的可能。', 0.5, 1.2, 6.5, 0.8, 13)
Stat(s, '互证结论：三条独立路径收敛至同一流形，max Δ < 3%', 0.5, 2.2, 6.5)
Stat(s, '≤10% 中断率约束下，关键 Knee 面积仅 1.646 m²', 0.5, 2.7, 6.5)
Img(s, 'python_optimization/algo_comparison.png', 7.2, 1.2, 5.5, 3.5)
Img(s, 'python_optimization/viz_wall_canvas.png', 0.5, 3.3, 2.8, 1.8)
Img(s, 'python_optimization/viz_parallel_coords.png', 3.5, 3.3, 3.2, 1.8)
T(s, 'Panichella, A. (2019). An adaptive evolutionary algorithm based on non-Euclidean geometry for many-objective optimization. GECCO 2019.', 0.5, 5.5, 12, 0.4, 9, False, MU)

# ═══════════ SLIDE 4 ═══════════
s = blank()
H(s, '三、帕累托解集形态特征与拓扑突变简析')
T(s, 'Knee 解规律：xc≈5.26, zc≈1.29 精确收敛于衍射+KDE质心。', 0.5, 1.2, 5.5, 0.5, 13)
T(s, 'Knee Jump (~9.7%) — 正交轴向転置\n  Lz 0.54→0.20m (−63%), 触及衍射下限;\n  Lx 7.32→9.31m (+27%), 水平补偿。', 0.5, 1.8, 5.5, 1.0, 12)
T(s, 'Jump B (~13.0%) — 单轴极限压缩\n  Lz 锁定, Lx 2.22→0.15m (−93%),\n  窗口退化为狭缝波导机构。', 0.5, 3.0, 5.5, 1.0, 12)
Img(s, 'python_optimization/jump_wall_canvas.png', 6.5, 1.2, 6, 5.5)

# ═══════════ SLIDE 5 ═══════════
s = blank()
H(s, '四、BO-LGBM 物理代理模型构建与高精度流形验证')
T(s, '混合数据集：5,000 LHS + 300 真实 Pareto + 2,000 边界过采样 = 7,300 样本\nLightGBM（树分裂天然适配 sinc 非连续响应面）\nOptuna 50 轮 / 5-Fold CV / feasible-MAE 为目标\n交叉切片验证：以 Knee 为锚, 4D 各扫 60 步, 模型与物理紧密贴合', 0.5, 1.2, 6.5, 2.2, 12)
Img(s, 'surrogate_lgbm/validation_sweeps.png', 7.5, 1.2, 5.2, 3.5)
Tbl(s, [['指标','值','阈值'],['Feasible MAE','0.52%','—'],['HV Deviation','1.19%','< 5%'],['IGD','0.0120','< 0.05']], 0.5, 3.8, [2.2,1.8,1.5])

# ═══════════ SLIDE 6 ═══════════
s = blank()
H(s, '五、两阶段 Warm-Start 混合最优化框架设计')
T(s, '代理悖论：全局 MAE 0.5%, 但树模型在可行域边界过光滑 → 局部系统性偏差。\n纯代理 Knee 1.755 m², 偏离物理真值 1.646 m² 达 6.7%。', 0.5, 1.2, 12, 0.8, 14)
Card(s, 'Phase 1 — 代理粗筛', 'NSGA-II 在 LGBM 上 200 代 (12s)。剪除 90% 不可行空间, 输出 300 个近似 Pareto 候选。', 1.5, 2.5, 4.5, 1.2)
T(s, '↓', 6.3, 2.8, 0.8, 0.5, 24, True, C2, PP_ALIGN.CENTER)
Card(s, 'Phase 2 — 物理精炼', '250 精英 + 50 随机混合注入（防止种群退化），物理 AGE-MOEA / NSGA-II 仅跑 20 代在线精炼，校正局部偏误。', 7.3, 2.5, 4.5, 1.2, True)
T(s, 'Shen, Y. & Pan, Y. (2023). BIM-supported energy performance analysis using explainable ML and multi-objective optimization. Applied Energy, 333, 120575.\nPanichella, A. (2022). Improved Pareto front modeling for large-scale many-objective optimization. GECCO 2022, 565–573.', 0.5, 4.3, 12, 0.8, 9, False, MU)

# ═══════════ SLIDE 7 ═══════════
s = blank()
H(s, '六、消融分析：基因多样性消融与边界约束')
T(s, '基因多样性消融：纯物理迭代 (深蓝线) 缓慢下降收敛；Warm-Start 管线 (黄、绿线) 从 gen 0 起 σ 即收缩至极小值——AI 先验知识在演化启动前已剪除冗余搜索空间。', 0.5, 1.2, 7, 1.0, 12)
Img(s, 'surrogate_lgbm/angle1_diversity.png', 8, 1.0, 5, 2.8)
T(s, '约束违规率拦截：纯物理 NSGA-II 早期违规率高达 60–70%，大量交叉子代越界造成无效算力损耗。Warm-Start 自 gen 0 起违规率趋近于 0——LGBM 离线已将硬性边界约束固化为不可违反的判定规则。', 0.5, 3.2, 7, 1.0, 12)
Img(s, 'surrogate_lgbm/angle3_violations.png', 8, 3.9, 5, 2.8)

# ═══════════ SLIDE 8 ═══════════
s = blank()
H(s, '七、四条技术路线 Pareto Front 比较与场景抉择')
Tbl(s, [['Arm','时间','Knee'],['Pure Physics','217s','1.646 m²'],['Pure Surrogate','12s','1.755 m²'],['Dual NSGA-II','45s','1.692 m²'],['Dual AGE-MOEA','84s','1.658 m²']], 0.5, 1.2, [3,1.5,2])
T(s, '深蓝（物理真值）、红（纯代理）、绿（Dual NSGA2）、橙（Dual AGE）。\nDual AGE-MOEA 面积仅偏离 ground truth 0.7%, 物理评估次数削减 90%。\n纯代理 (12s) 和 Dual NSGA-II (45s) 为算力受限场景下的高性价比备选。', 0.5, 3.2, 6.5, 1.2, 12)
Img(s, 'surrogate_lgbm/angle4_pareto.png', 7.5, 1.2, 5.2, 5.5)

# ═══════════ SLIDE 9 ═══════════
s = blank()
H(s, '八、全链路工作流闭环总复盘')
steps = [
    ('1','物理网格化降噪 (→S2)：消除轨迹随机抖动，R²=0.9941'),
    ('2','三算法机制互证 (→S3)：独立路径收敛至同一前沿，max Δ<3%'),
    ('3','拓扑特征解耦 (→S4)：Knee Jump (正交転置) + Jump B (单轴压缩)'),
    ('4','代理平替与流形验证 (→S5)：R²=0.995, IGD=0.0120'),
    ('5','两阶段混合管线落地 (→S6-8)：12s 离线剪枝 + 20 代精炼，90% 算力节省, 0.7% 逼近'),
]
for i,(num,desc) in enumerate(steps):
    y = 1.2 + i*1.1
    # Circle
    shp = s.shapes.add_shape(9, Inches(0.8), Inches(y), Inches(0.5), Inches(0.5))  # oval
    shp.fill.solid(); shp.fill.fore_color.rgb = C1
    shp.line.fill.background()
    shp.text_frame.paragraphs[0].text = num; shp.text_frame.paragraphs[0].font.size = Pt(14)
    shp.text_frame.paragraphs[0].font.bold = True; shp.text_frame.paragraphs[0].font.color.rgb = WH
    shp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    T(s, desc, 1.5, y, 11, 0.8, 13)

# ═══════════ SLIDE 10 ═══════════
s = blank()
H(s, '九、多路线试错负结果与全局收敛边界的确证')
Tbl(s, [['方法','失败原因'],['RBF 局部精炼','30/300/500 点三版迭代无一改进'],
          ['置信域序列优化','3 轮 FAILURE, 半径砍至 1.25%'],
          ['Ridge+LGBM Stacking','全局 MAE 优化但 Knee 劣化 1.66→1.73'],
          ['GP+LGBM Stacking','O(n³) 算力限制, feasible MAE 退化至 1.04%']], 0.5, 1.2, [3.5,7])
# Conclusion box
shp = s.shapes.add_shape(5, Inches(0.8), Inches(3.5), Inches(11.5), Inches(2.5))
shp.fill.solid(); shp.fill.fore_color.rgb = C1; shp.line.fill.background()
tf = shp.text_frame; tf.word_wrap = True; tf.margin_left = Pt(20); tf.margin_right = Pt(20)
p = tf.paragraphs[0]; p.text = '四项进阶技术路线均未突破当前 Pipeline 的最优解——这并非方法缺陷，而是物理空间本身的约束极限。'; p.font.size = Pt(14); p.font.color.rgb = WH; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = '\n负结果从反面构成了逻辑闭环：\nBO-LGBM 两阶段 Warm-Start 混合管线已探明该电磁约束下的全局收敛上界。'; p2.font.size = Pt(14); p2.font.color.rgb = WH; p2.alignment = PP_ALIGN.CENTER

# Save
prs.save('slides.pptx')
print('✅ slides.pptx saved')
